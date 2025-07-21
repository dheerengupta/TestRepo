import os
import sys
import argparse

import openai
from pptx import Presentation
from pptx.util import Inches, Pt


def generate_title_and_subtitle(topic: str) -> tuple[str, str]:
    """Use OpenAI to craft a professional title and subtitle for the presentation."""
    # Ensure the API key is available
    if not os.getenv("OPENAI_API_KEY"):
        raise EnvironmentError("OPENAI_API_KEY environment variable is not set.")

    openai.api_key = os.environ["OPENAI_API_KEY"]

    # Craft the prompt
    prompt = (
        "You are a presentation expert. Generate a compelling, professional PowerPoint "
        "title and subtitle for a presentation whose subject is: \"{}\".\n".format(topic)
    )
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo-16k",  # Use an economical but capable model
        messages=[{"role": "user", "content": prompt}],
        max_tokens=60,
    )

    content = response.choices[0].message.content.strip()
    # Expecting output in the format: Title: ...\nSubtitle: ...
    title = None
    subtitle = None
    for line in content.splitlines():
        if line.lower().startswith("title:"):
            title = line.split(":", 1)[1].strip()
        elif line.lower().startswith("subtitle:"):
            subtitle = line.split(":", 1)[1].strip()
    # Fallback: treat first line as title, second as subtitle
    if not title:
        lines = content.splitlines()
        title = lines[0].strip()
        subtitle = lines[1].strip() if len(lines) > 1 else ""
    return title, subtitle


def create_pptx(title: str, subtitle: str, output_path: str) -> None:
    """Create a PowerPoint file with a single title slide."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]  # Usually the title slide
    slide = prs.slides.add_slide(title_slide_layout)

    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    # Set title
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    # Basic formatting tweaks
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(40)
    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)

    prs.save(output_path)


def main():
    parser = argparse.ArgumentParser(description="Generate the first slide of a PowerPoint presentation using AI.")
    parser.add_argument("topic", help="Topic of the presentation")
    parser.add_argument("-o", "--output", default="presentation.pptx", help="Output PPTX file path")
    args = parser.parse_args()

    print("Generating title and subtitle using OpenAI ...")
    title, subtitle = generate_title_and_subtitle(args.topic)
    print(f"Title: {title}\nSubtitle: {subtitle}")

    print(f"Creating PowerPoint file at {args.output} ...")
    create_pptx(title, subtitle, args.output)
    print("Done! Your presentation has been saved.")


if __name__ == "__main__":
    main()